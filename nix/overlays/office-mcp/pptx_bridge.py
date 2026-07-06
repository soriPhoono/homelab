#!/usr/bin/env python3
"""Bridge script for python-pptx operations.

Called by the office-mcp JS server as:
  python3 pptx_bridge.py <command> '<json_args>'

Outputs JSON to stdout.
"""
import sys
import json
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from lxml import etree


def read_structure(args):
    """Read detailed shape/text/table info per slide."""
    prs = Presentation(args["path"])
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "layout": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": [],
        }

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    runs = []
                    for run in para.runs:
                        runs.append({
                            "text": run.text,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "font_size": (
                                str(run.font.size) if run.font.size else None
                            ),
                            "font_name": run.font.name,
                            "color": (
                                str(run.font.color.rgb)
                                if run.font.color and run.font.color.rgb
                                else None
                            ),
                        })
                    paragraphs.append({
                        "text": para.text,
                        "level": para.level,
                        "runs": runs,
                    })
                shape_info["text_frame"] = paragraphs

            if shape.has_table:
                table_data = {
                    "rows": [],
                    "col_count": shape.table.columns.__len__(),
                }
                for row in shape.table.rows:
                    table_data["rows"].append([
                        cell.text for cell in row.cells
                    ])
                shape_info["table"] = table_data

            slide_data["shapes"].append(shape_info)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                slide_data["notes"] = notes

        slides.append(slide_data)

    return {
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides,
    }


def replace_text(args):
    """Find/replace text across all slides, preserving formatting."""
    prs = Presentation(args["path"])
    replacements = args["replacements"]
    count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for old, new in replacements.items():
                        if old not in para.text:
                            continue
                        # Run-level replacement
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1
                        # Cross-run fallback
                        if old in para.text:
                            full = para.text.replace(old, new)
                            if para.runs:
                                para.runs[0].text = full
                                for r in list(para.runs[1:]):
                                    r.text = ""
                                count += 1

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for old, new in replacements.items():
                                if old not in para.text:
                                    continue
                                for run in para.runs:
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                        count += 1
                                if old in para.text:
                                    full = para.text.replace(old, new)
                                    if para.runs:
                                        para.runs[0].text = full
                                        for r in list(para.runs[1:]):
                                            r.text = ""
                                        count += 1

    output = args.get("output", args["path"])
    prs.save(output)
    return {"success": True, "replacements_made": count, "output": output}


def add_slide(args):
    """Add a slide using python-pptx, preserving existing slides."""
    prs = Presentation(args["path"])
    layout_index = args.get("layoutIndex", 1)
    content = args.get("content", {})

    if layout_index >= len(prs.slide_layouts):
        layout_index = len(prs.slide_layouts) - 1

    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    # Set title
    if "title" in content and slide.shapes.title:
        slide.shapes.title.text = content["title"]

    # Set body text
    if "body" in content:
        body = content["body"]
        if isinstance(body, str):
            body = [body]
        # Find body placeholder
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            for i, line in enumerate(body):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

    # Notes
    if "notes" in content:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content["notes"]

    output = args.get("output", args["path"])
    prs.save(output)
    return {
        "success": True,
        "slide_number": len(prs.slides),
        "total_slides": len(prs.slides),
        "output": output,
    }


def delete_slide(args):
    """Delete a slide by number (1-based), preserving all others."""
    prs = Presentation(args["path"])
    slide_number = args["slideNumber"]

    if slide_number < 1 or slide_number > len(prs.slides):
        return {
            "error": f"Slide {slide_number} out of range (1-{len(prs.slides)})",
        }

    # Access the internal XML to remove the slide
    rId = prs.slides._sldIdLst[slide_number - 1].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)

    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[slide_number - 1]
    sldIdLst.remove(sldId)

    output = args.get("output", args["path"])
    prs.save(output)
    return {
        "success": True,
        "deleted_slide": slide_number,
        "remaining_slides": len(prs.slides),
        "output": output,
    }


def modify_slide_text(args):
    """Replace text on a specific slide, preserving formatting."""
    prs = Presentation(args["path"])
    slide_number = args["slideNumber"]
    replacements = args["replacements"]

    if slide_number < 1 or slide_number > len(prs.slides):
        return {
            "error": f"Slide {slide_number} out of range (1-{len(prs.slides)})",
        }

    slide = prs.slides[slide_number - 1]
    count = 0

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for old, new in replacements.items():
                    if old not in para.text:
                        continue
                    for run in para.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            count += 1
                    if old in para.text:
                        full = para.text.replace(old, new)
                        if para.runs:
                            para.runs[0].text = full
                            for r in list(para.runs[1:]):
                                r.text = ""
                            count += 1

    output = args.get("output", args["path"])
    prs.save(output)
    return {
        "success": True,
        "replacements_made": count,
        "slide": slide_number,
        "output": output,
    }


def update_table_cell(args):
    """Update a specific table cell on a slide."""
    prs = Presentation(args["path"])
    slide_number = args["slideNumber"]
    table_index = args.get("tableIndex", 0)
    row = args["row"]
    col = args["col"]
    value = args["value"]

    if slide_number < 1 or slide_number > len(prs.slides):
        return {
            "error": f"Slide {slide_number} out of range (1-{len(prs.slides)})",
        }

    slide = prs.slides[slide_number - 1]

    # Find table shapes
    tables = [s for s in slide.shapes if s.has_table]
    if table_index >= len(tables):
        return {
            "error": f"Table index {table_index} out of range (0-{len(tables) - 1})",
        }

    table = tables[table_index].table
    if row >= len(table.rows) or col >= len(table.columns):
        return {
            "error": (
                f"Cell ({row},{col}) out of range "
                f"(rows: {len(table.rows)}, cols: {len(table.columns)})"
            ),
        }

    cell = table.rows[row].cells[col]
    # Preserve existing formatting by updating run text
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].text = str(value)
        for r in list(cell.text_frame.paragraphs[0].runs[1:]):
            r.text = ""
    else:
        cell.text = str(value)

    output = args.get("output", args["path"])
    prs.save(output)
    return {
        "success": True,
        "slide": slide_number,
        "table": table_index,
        "row": row,
        "col": col,
        "value": value,
        "output": output,
    }


def duplicate_slide(args):
    """Duplicate a slide (deep copy with all formatting)."""
    prs = Presentation(args["path"])
    slide_number = args["slideNumber"]

    if slide_number < 1 or slide_number > len(prs.slides):
        return {
            "error": f"Slide {slide_number} out of range (1-{len(prs.slides)})",
        }

    # Use lxml deep copy of the slide XML
    template_slide = prs.slides[slide_number - 1]
    slide_layout = template_slide.slide_layout

    # Add a new slide with the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all shape XML from template to new slide
    # Clear the new slide's default shapes first
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Deep copy shapes from template
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Copy slide relationships (images, etc.)
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    # Copy notes
    if template_slide.has_notes_slide:
        notes_text = template_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            new_slide.notes_slide.notes_text_frame.text = notes_text

    output = args.get("output", args["path"])
    prs.save(output)
    return {
        "success": True,
        "original_slide": slide_number,
        "duplicated_as": len(prs.slides),
        "total_slides": len(prs.slides),
        "output": output,
    }


COMMANDS = {
    "read_structure": read_structure,
    "replace_text": replace_text,
    "add_slide": add_slide,
    "delete_slide": delete_slide,
    "modify_slide_text": modify_slide_text,
    "update_table_cell": update_table_cell,
    "duplicate_slide": duplicate_slide,
}


def main():
    if len(sys.argv) < 3:
        print(json.dumps({"error": "Usage: pptx_bridge.py <command> '<json_args>'"}))
        sys.exit(1)

    cmd = sys.argv[1]
    try:
        cmd_args = json.loads(sys.argv[2])
    except json.JSONDecodeError as e:
        print(json.dumps({"error": f"Invalid JSON args: {e}"}))
        sys.exit(1)

    if cmd not in COMMANDS:
        print(json.dumps({"error": f"Unknown command: {cmd}"}))
        sys.exit(1)

    try:
        result = COMMANDS[cmd](cmd_args)
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({"error": str(e), "command": cmd}))
        sys.exit(1)


if __name__ == "__main__":
    main()
